import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# 1. Initialize Presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 2. Color Palette Definitions
BG_COLOR = RGBColor(18, 24, 38)          # #121826 (Very Dark Navy)
CARD_BG = RGBColor(30, 41, 59)           # #1E293B (Dark Slate/Grey-Blue)
TEXT_PRIMARY = RGBColor(248, 250, 252)   # #F8FAFC (Off-white)
TEXT_SECONDARY = RGBColor(148, 163, 184) # #94A3B8 (Cool Grey)
ACCENT_TEAL = RGBColor(34, 211, 238)     # #22D3EE (Bright Teal/Cyan)
ACCENT_INDIGO = RGBColor(129, 140, 248)   # #818CF8 (Bright Indigo/Purple)
ACCENT_ORANGE = RGBColor(251, 146, 60)   # #FB923C (Warm Orange for emphasis/alerts)
BORDER_COLOR = RGBColor(51, 65, 85)       # #334155 (Medium Slate)

# 3. Helper Functions
def apply_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_header(slide, section_tag, slide_title):
    # Section Tag (small text at top)
    tag_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.3))
    tf_tag = tag_box.text_frame
    tf_tag.word_wrap = True
    p_tag = tf_tag.paragraphs[0]
    p_tag.text = section_tag.upper()
    p_tag.font.name = "Apple SD Gothic Neo"
    p_tag.font.size = Pt(10)
    p_tag.font.bold = True
    p_tag.font.color.rgb = ACCENT_INDIGO
    
    # Slide Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.65), Inches(11.733), Inches(0.6))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = slide_title
    p_title.font.name = "Apple SD Gothic Neo"
    p_title.font.size = Pt(24)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_PRIMARY

def draw_card(slide, left, top, width, height, title_text="", bg_color=CARD_BG, border_color=BORDER_COLOR):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)
    
    if title_text:
        txBox = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.15), width - Inches(0.3), Inches(0.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Apple SD Gothic Neo"
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = ACCENT_TEAL
        return left + Inches(0.25), top + Inches(0.65), width - Inches(0.5), height - Inches(0.8)
        
    return left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4)

def add_bullet(tf, prefix, text, prefix_color=ACCENT_TEAL):
    p = tf.add_paragraph()
    p.space_before = Pt(6)
    p.space_after = Pt(6)
    p.font.name = "Apple SD Gothic Neo"
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_PRIMARY
    
    # Bullet/Prefix
    r_pref = p.add_run()
    r_pref.text = "• " + prefix
    r_pref.font.bold = True
    r_pref.font.color.rgb = prefix_color
    
    # Body Text
    r_text = p.add_run()
    r_text.text = text
    r_text.font.color.rgb = TEXT_PRIMARY

# ==========================================================
# SLIDE 1: Title Slide
# ==========================================================
slide_layout = prs.slide_layouts[6] # Blank
slide = prs.slides.add_slide(slide_layout)
apply_slide_background(slide)

# Decorative left stripe
stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
stripe.fill.solid()
stripe.fill.fore_color.rgb = ACCENT_TEAL
stripe.line.fill.background()

# Subtitle Header tag
tag_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(0.4))
tf_tag = tag_box.text_frame
p_tag = tf_tag.paragraphs[0]
p_tag.text = "MASTER'S DEGREE THESIS PRESENTATION"
p_tag.font.name = "Apple SD Gothic Neo"
p_tag.font.size = Pt(13)
p_tag.font.bold = True
p_tag.font.color.rgb = ACCENT_INDIGO

# Title Box
title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(2.2))
tf_title = title_box.text_frame
tf_title.word_wrap = True
p_title = tf_title.paragraphs[0]
p_title.text = "zk-SNARK 기반, 프라이버시를 보존하는\n수치형 의료 데이터에 대한 증명 검색 시스템"
p_title.font.name = "Apple SD Gothic Neo"
p_title.font.size = Pt(34)
p_title.font.bold = True
p_title.font.color.rgb = TEXT_PRIMARY

# English Title
p_title_en = tf_title.add_paragraph()
p_title_en.text = "Privacy-Preserving Proof Search System for Numerical Medical Data based on zk-SNARK"
p_title_en.font.name = "Apple SD Gothic Neo"
p_title_en.font.size = Pt(17)
p_title_en.font.color.rgb = ACCENT_TEAL
p_title_en.space_before = Pt(12)

# Author info
info_box = slide.shapes.add_textbox(Inches(1.0), Inches(5.0), Inches(11.3), Inches(1.5))
tf_info = info_box.text_frame
p_info1 = tf_info.paragraphs[0]
p_info1.text = "서강대학교 대학원 컴퓨터공학과"
p_info1.font.name = "Apple SD Gothic Neo"
p_info1.font.size = Pt(14)
p_info1.font.color.rgb = TEXT_SECONDARY

p_info2 = tf_info.add_paragraph()
p_info2.text = "발표자 : 황 보 진 우 (지도교수 박수용)"
p_info2.font.name = "Apple SD Gothic Neo"
p_info2.font.size = Pt(15)
p_info2.font.bold = True
p_info2.font.color.rgb = TEXT_PRIMARY
p_info2.space_before = Pt(6)

p_info3 = tf_info.add_paragraph()
p_info3.text = "2026년 6월"
p_info3.font.name = "Apple SD Gothic Neo"
p_info3.font.size = Pt(12)
p_info3.font.color.rgb = TEXT_SECONDARY
p_info3.space_before = Pt(6)

# ==========================================================
# SLIDE 2: 연구 배경 (Research Background)
# ==========================================================
slide = prs.slides.add_slide(slide_layout)
apply_slide_background(slide)
add_header(slide, "01. Introduction", "연구 배경: 의료 데이터 유출 및 최소화 원칙")

# 3 Columns
w = Inches(3.6)
g = Inches(0.4)
top = Inches(1.6)
height = Inches(4.8)

# Col 1: 유출 피해
il, it, iw, ih = draw_card(slide, Inches(0.8), top, w, height, "의료 데이터 유출 피해의 심각성")
tx = slide.shapes.add_textbox(il, it, iw, ih)
tf = tx.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "■ 민감한 데이터 가치 & 보안 요구"
p.font.name = "Apple SD Gothic Neo"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = ACCENT_ORANGE

add_bullet(tf, "피해 비용 1위: ", "의료 산업은 14년 연속 전 세계 데이터 유출 평균 피해 비용이 가장 높은 분야 (IBM 2025 보고서)")
add_bullet(tf, "평균 742만 달러: ", "침해로 인한 사고당 평균 비용이 742만 달러(USD)에 육박")
add_bullet(tf, "느린 대응 속도: ", "유출 사고의 식별 및 통제에 평균 279일 소요되어 피해 누적 극대화")

# Col 2: 데이터 최소화
il, it, iw, ih = draw_card(slide, Inches(0.8) + w + g, top, w, height, "데이터 최소화 (Data Minimization) 원칙")
tx = slide.shapes.add_textbox(il, it, iw, ih)
tf = tx.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "■ 프라이버시 침해의 근본적 예방"
p.font.name = "Apple SD Gothic Neo"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = ACCENT_INDIGO

add_bullet(tf, "필수 정보만 수집: ", "시스템 운영에 반드시 필요한 데이터만 수집/보존하고 외부는 배제 (GDPR 권고사항)")
add_bullet(tf, "보안 강화의 핵심: ", "\"수집되지 않은 데이터는 사람들에게 피해를 줄 수 없다\"는 철학 기반")
add_bullet(tf, "사용자 친화적 설계: ", "데이터 관리 범위 자체를 축소하여 보안 취약점을 최소화하는 확실한 접근법")

# Col 3: ZKP 도입
il, it, iw, ih = draw_card(slide, Inches(0.8) + (w + g)*2, top, w, height, "영지식 증명 (Zero-Knowledge Proof)")
tx = slide.shapes.add_textbox(il, it, iw, ih)
tf = tx.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "■ 수학적 프라이버시 보존 수단"
p.font.name = "Apple SD Gothic Neo"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = ACCENT_TEAL

add_bullet(tf, "원문 노출 제로: ", "개인의 민감한 건강 수치(원문)를 타인이나 서버에 노출하지 않음")
add_bullet(tf, "수학적 팩트 검증: ", "원본을 감춘 상태에서 해당 데이터가 '특정 조건을 만족한다'는 사실만을 입증 가능")
add_bullet(tf, "의료 데이터 유통망: ", "개인정보 침해 없는 안전한 데이터 거래 및 유통 플랫폼 구축을 위한 핵심 암호학 기술")


# ==========================================================
# SLIDE 3: 기존 연구 및 한계점 (Existing Research & Limitations)
# ==========================================================
slide = prs.slides.add_slide(slide_layout)
apply_slide_background(slide)
add_header(slide, "01. Introduction", "기존 연구의 흐름 및 한계점")

# 2 Columns (Left: General Architecture, Right: Limitations)
w2 = Inches(5.65)
g2 = Inches(0.4)

il, it, iw, ih = draw_card(slide, Inches(0.8), top, w2, height, "일반적인 영지식 기반 데이터 거래 (Gao et al. 2024)")
tx = slide.shapes.add_textbox(il, it, iw, ih)
tf = tx.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "■ 분산 저장소 기반 일회성 ZKP 생성 아키텍처"
p.font.name = "Apple SD Gothic Neo"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = ACCENT_INDIGO

add_bullet(tf, "데이터 위탁 저장: ", "소유자가 데이터를 IPFS(분산 파일 시스템) 등 제3자 저장소에 암호화하여 저장")
add_bullet(tf, "블록체인 연계: ", "암호화 데이터의 해시를 블록체인에 등록하여 무결성 검증에 사용")
add_bullet(tf, "실시간 생성 & 검증: ", "수요 발생 시 소유자가 실시간으로 영지식 증명을 생성하고 스마트 컨트랙트로 검증")
add_bullet(tf, "일회성 폐기: ", "검증된 증명은 일시적인 검사 후 즉시 폐기되는 프로세스가 일반적임")

il, it, iw, ih = draw_card(slide, Inches(0.8) + w2 + g2, top, w2, height, "기존 탈중앙화 데이터 시스템의 치명적 한계")
tx = slide.shapes.add_textbox(il, it, iw, ih)
tf = tx.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "■ 실용성과 확장성 저해 요인"
p.font.name = "Apple SD Gothic Neo"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = ACCENT_ORANGE

add_bullet(tf, "데이터 주권 약화: ", "암호화했다 하더라도 결국 데이터 원문을 제3자 네트워크(IPFS 등)에 영구 위탁해야 함")
add_bullet(tf, "확장성 결여 (1:1 통신): ", "소유자와 수요자 간의 개별적인 실시간 증명 프로토콜로 인해 다수 참여 시 지연 가속화")
add_bullet(tf, "대규모 검색 불가: ", "수요자가 원하는 데이터의 유무조차 실시간으로 파악하기 힘들어 대용량 데이터 유통 플랫폼에 부적합")


# ==========================================================
# SLIDE 4: 문제 정의 (Problem Definition)
# ==========================================================
slide = prs.slides.add_slide(slide_layout)
apply_slide_background(slide)
add_header(slide, "01. Introduction", "문제 정의: 프라이버시와 검색 가능성의 이중성")

# Center Big Problem Question Card
q_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.733), Inches(1.2))
tf_q = q_box.text_frame
tf_q.word_wrap = True
p_q = tf_q.paragraphs[0]
p_q.text = "핵심 연구 문제 (Core Research Question)"
p_q.font.name = "Apple SD Gothic Neo"
p_q.font.size = Pt(14)
p_q.font.bold = True
p_q.font.color.rgb = ACCENT_INDIGO

p_q2 = tf_q.add_paragraph()
p_q2.text = "\"어떻게 하면 데이터 원문을 은닉(Privacy)하여 데이터 최소화를 실현하면서도,\n 실시간 검색(Searchability)이 가능한 프라이버시 검색 시스템을 구축할 것인가?\""
p_q2.font.name = "Apple SD Gothic Neo"
p_q2.font.size = Pt(19)
p_q2.font.bold = True
p_q2.font.color.rgb = ACCENT_TEAL
p_q2.space_before = Pt(8)

# 2 Columns comparing approach 1 and 2
top_sub = Inches(2.9)
h_sub = Inches(3.8)

il, it, iw, ih = draw_card(slide, Inches(0.8), top_sub, w2, h_sub, "접근 ① : 정보 노출형 검색 (보안성 결여)")
tx = slide.shapes.add_textbox(il, it, iw, ih)
tf = tx.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "■ 순서 보존 암호화 (OPE) 계열 (예: CryptDB)"
p.font.name = "Apple SD Gothic Neo"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = ACCENT_ORANGE

add_bullet(tf, "작동 방식: ", "암호문 상태에서 데이터의 대소 관계와 순서를 보존하여 범위 검색 수행")
add_bullet(tf, "치명적 한계: ", "데이터 분포 정보가 노출되므로 추론 공격에 치명적임")
add_bullet(tf, "취약점 사례: ", "Naveed et al. (2015) 실험 결과, 실제 OPE 데이터베이스 공격 시 95% 이상의 환자 원문 기록 복구 성공")

il, it, iw, ih = draw_card(slide, Inches(0.8) + w2 + g2, top_sub, w2, h_sub, "접근 ② : 완전 은닉형 검색 (실용성 결여)")
tx = slide.shapes.add_textbox(il, it, iw, ih)
tf = tx.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "■ 검증 가능한 DB (Verifiable DB) 계열 (예: ZKSQL, IntegriDB)"
p.font.name = "Apple SD Gothic Neo"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = ACCENT_ORANGE

add_bullet(tf, "작동 방식: ", "쿼리 요청 시마다 실시간으로 영지식 서킷 연산을 수행하여 결과 무결성 검증")
add_bullet(tf, "치명적 한계: ", "수 밀리초(ms) 수준의 일반 SQL 질의가 ZKSQL 하에선 수 초 ~ 수 분 소요")
add_bullet(tf, "운영 오버헤드: ", "대화형 프로토콜로 인한 통신 지연 발생, 실시간 변화하는 동적 DB에 적용 불가")


# ==========================================================
# SLIDE 5: 관련 연구 분류 및 한계점 (Related Research & Limitations)
# ==========================================================
slide = prs.slides.add_slide(slide_layout)
apply_slide_background(slide)
add_header(slide, "02. Related Research", "관련 연구 분류 및 한계점 비교")

# 4 Columns
w4 = Inches(2.74)
g4 = Inches(0.25)
top4 = Inches(1.6)
height4 = Inches(4.8)

# Col 1: 암호화 검색 및 OPE
il1, it1, iw1, ih1 = draw_card(slide, Inches(0.8), top4, w4, height4, "암호화 검색 및 OPE")
tx1 = slide.shapes.add_textbox(il1, it1, iw1, ih1)
tf1 = tx1.text_frame
tf1.word_wrap = True
p1 = tf1.paragraphs[0]
p1.text = "■ 순서 보존 암호 (OPE)"
p1.font.name = "Apple SD Gothic Neo"
p1.font.size = Pt(13)
p1.font.bold = True
p1.font.color.rgb = ACCENT_INDIGO
add_bullet(tf1, "대표 기법: ", "CryptDB, OPE (Agrawal et al.) 등")
add_bullet(tf1, "핵심 특징: ", "암호문 상태에서 대소 관계 및 순서 보존 검색 지원")
add_bullet(tf1, "치명적 한계: ", "순서/빈도 정보 노출로 추론 공격에 취약, 원문 복구 위협 존재")

# Col 2: 동형 암호 및 인덱스 검색
il2, it2, iw2, ih2 = draw_card(slide, Inches(0.8) + w4 + g4, top4, w4, height4, "동형 암호 & 인덱스 검색")
tx2 = slide.shapes.add_textbox(il2, it2, iw2, ih2)
tf2 = tx2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.text = "■ 부분 동형 암호 & 인덱스"
p2.font.name = "Apple SD Gothic Neo"
p2.font.size = Pt(13)
p2.font.bold = True
p2.font.color.rgb = ACCENT_INDIGO
add_bullet(tf2, "대표 기법: ", "ASPE (Wang et al.), Paillier 결합 머신러닝 인덱스 등")
add_bullet(tf2, "핵심 특징: ", "동형 암호 연산 및 다차원 인덱스(R-tree) 구조 활용")
add_bullet(tf2, "치명적 한계: ", "서버의 복호화 키 보유 또는 중간 결과 복호화 필요 (SPOF)")

# Col 3: 검증 가능한 데이터베이스
il3, it3, iw3, ih3 = draw_card(slide, Inches(0.8) + (w4 + g4)*2, top4, w4, height4, "검증 가능한 DB")
tx3 = slide.shapes.add_textbox(il3, it3, iw3, ih3)
tf3 = tx3.text_frame
tf3.word_wrap = True
p3 = tf3.paragraphs[0]
p3.text = "■ 쿼리 무결성 증명 (VDB)"
p3.font.name = "Apple SD Gothic Neo"
p3.font.size = Pt(13)
p3.font.bold = True
p3.font.color.rgb = ACCENT_INDIGO
add_bullet(tf3, "대표 기법: ", "IntegriDB, vSQL (Zhang et al.), ZKSQL 등")
add_bullet(tf3, "핵심 특징: ", "질의 결과의 정확성과 완전성(Sound & Complete)을 암호학적으로 증명")
add_bullet(tf3, "치명적 한계: ", "실시간 서킷 증명 생성 비용 과다 (수 초~수 분 소요, 실시간성 부족)")

# Col 4: 하드웨어 및 블록체인
il4, it4, iw4, ih4 = draw_card(slide, Inches(0.8) + (w4 + g4)*3, top4, w4, height4, "하드웨어 & 블록체인")
tx4 = slide.shapes.add_textbox(il4, it4, iw4, ih4)
tf4 = tx4.text_frame
tf4.word_wrap = True
p4 = tf4.paragraphs[0]
p4.text = "■ TEE 및 블록체인 활용"
p4.font.name = "Apple SD Gothic Neo"
p4.font.size = Pt(13)
p4.font.bold = True
p4.font.color.rgb = ACCENT_INDIGO
add_bullet(tf4, "대표 기법: ", "VeriDB (Intel SGX), vChain (Blockchain ADS) 등")
add_bullet(tf4, "핵심 특징: ", "신뢰 실행 환경 내 연산 처리 또는 분산 원장 기반 무결성 검증")
add_bullet(tf4, "치명적 한계: ", "SGX 부채널 공격 취약/제조사 의존, 블록체인의 기밀성 결여")


# ==========================================================
# SLIDE 6: 제안 시스템 핵심 아이디어 (Proposed Ideas)
# ==========================================================
slide = prs.slides.add_slide(slide_layout)
apply_slide_background(slide)
add_header(slide, "03. Proposed System", "제안 시스템 핵심 아이디어")

# 3 Columns
il, it, iw, ih = draw_card(slide, Inches(0.8), top, w, height, "① 원문 미저장 증명 검색 (Proof Search)")
tx = slide.shapes.add_textbox(il, it, iw, ih)
tf = tx.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "■ 원문 보관 없는 검색 인프라"
p.font.name = "Apple SD Gothic Neo"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = ACCENT_TEAL

add_bullet(tf, "증명 데이터 저장: ", "데이터 원문은 소유자(DP) 로컬 단말기에만 보관, 서버에는 특정 범위 존재 유무를 증명하는 범위 증명(ZKRP)만 저장")
add_bullet(tf, "원문 복원 원천 차단: ", "서버가 완전히 해킹당하더라도 데이터 원문 유출 우려 원천 해제")
add_bullet(tf, "데이터 유통망 활성화: ", "수요자에게 데이터의 신뢰 가능한 존재를 공시하여 안전한 거래 촉진")

il, it, iw, ih = draw_card(slide, Inches(0.8) + w + g, top, w, height, "② 참고치 기반 오프라인 증명 생성")
tx = slide.shapes.add_textbox(il, it, iw, ih)
tf = tx.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "■ ZK-오버헤드 극복 및 실시간 검색"
p.font.name = "Apple SD Gothic Neo"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = ACCENT_TEAL

add_bullet(tf, "의료 참고치 활용: ", "의료 데이터 도메인 고유 속성(WHO 임상 참고치 등) 기반 범위 사전 정의")
add_bullet(tf, "오프라인 사전 생성: ", "실시간 쿼리 발생 시점에 증명을 생성하지 않고, 데이터 발급 시점에 범위 증명을 미리 생성")
add_bullet(tf, "O(1)의 데이터 펫치: ", "검색 시점에는 ZK 연산 없이 단순 DB 인덱스 조회만 수행하여 응답속도 획기적 단축")

il, it, iw, ih = draw_card(slide, Inches(0.8) + (w + g)*2, top, w, height, "③ 단일 그룹 커밋을 통한 소유권 검증")
tx = slide.shapes.add_textbox(il, it, iw, ih)
tf = tx.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "■ 복합 쿼리의 소유권 연계 증명"
p.font.name = "Apple SD Gothic Neo"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = ACCENT_TEAL

add_bullet(tf, "다중 증명 그룹화: ", "여러 개의 서브 쿼리를 만족하는 서로 다른 범위 증명이 동일 소유자의 것임을 보장")
add_bullet(tf, "단일 커밋 결합: ", "증명 식별자 해시(prf_hash) + DP_ID + nonce를 묶어 하나의 커밋 생성")
add_bullet(tf, "단 1회 검증: ", "서브 쿼리 개수와 상관없이 단 1회의 ZK 검증만으로 복합 조건 및 소유권을 수학적으로 증명")


# ==========================================================
# SLIDE 7: 구성 요소 및 아키텍처 (Components & Flow)
# ==========================================================
slide = prs.slides.add_slide(slide_layout)
apply_slide_background(slide)
add_header(slide, "03. Proposed System", "시스템 구성요소 및 역할 모델")

# 4 Columns
w4 = Inches(2.6)
g4 = Inches(0.4)

il, it, iw, ih = draw_card(slide, Inches(0.8), top, w4, height, "1. 발급기관 (Data Issuer)")
tx = slide.shapes.add_textbox(il, it, iw, ih)
tf = tx.text_frame
tf.word_wrap = True
add_bullet(tf, "주체: ", "병원, 검진기관 등 신뢰받는 개인 데이터 생성처")
add_bullet(tf, "역할: ", "데이터 발급 및 디지털 서명 부여")
add_bullet(tf, "증명 대행: ", "강력한 컴퓨팅 자원을 활용해 환자를 대행하여 범위 증명(ZKRP)을 오프라인으로 사전 생성")

il, it, iw, ih = draw_card(slide, Inches(0.8) + w4 + g4, top, w4, height, "2. 소유자 (Data Provider)")
tx = slide.shapes.add_textbox(il, it, iw, ih)
tf = tx.text_frame
tf.word_wrap = True
add_bullet(tf, "주체: ", "의료 데이터를 발급받는 개별 환자(DP)")
add_bullet(tf, "역할: ", "원본 데이터는 개인 지갑 앱 등 안전한 개인 스토리지에만 자율 보관")
add_bullet(tf, "주권 확보: ", "발급기관이 ZKRP를 생성하여 검색 서버에 전송하는 과정에 대한 디지털 동의(서명) 수행")

il, it, iw, ih = draw_card(slide, Inches(0.8) + (w4 + g4)*2, top, w4, height, "3. 검색 서버 (Search Server)")
tx = slide.shapes.add_textbox(il, it, iw, ih)
tf = tx.text_frame
tf.word_wrap = True
add_bullet(tf, "역할: ", "발송된 ZKRP 및 메타데이터를 데이터베이스(PROOF, MERKLE)에 적재. (원문 절대 비보존)")
add_bullet(tf, "쿼리 처리: ", "수요자(DC)의 속성/범위 검색 처리")
add_bullet(tf, "소유권 증명: ", "다중 조건 검색 시, 동일 소유자 확인을 위한 소유권 커밋 및 ZK 증명 실시간 생성")

il, it, iw, ih = draw_card(slide, Inches(0.8) + (w4 + g4)*3, top, w4, height, "4. 수요자 (Data Consumer)")
tx = slide.shapes.add_textbox(il, it, iw, ih)
tf = tx.text_frame
tf.word_wrap = True
add_bullet(tf, "주체: ", "제약회사, 연구소, 보험회사 등 의료 데이터 구매 희망자")
add_bullet(tf, "역할: ", "검색 서버에 쿼리를 전송하여 적합한 데이터의 개수와 증명 셋을 수신")
add_bullet(tf, "무결성 검증: ", "수신한 ZKRP 및 소유권 증명을 로컬에서 독자 검증하여 존재 신뢰")


# ==========================================================
# SLIDE 8: 범위 증명 및 회로 구조 (Range Proof Circuit)
# ==========================================================
slide = prs.slides.add_slide(slide_layout)
apply_slide_background(slide)
add_header(slide, "03. Proposed System", "범위 증명(ZKRP) 회로 및 검증 로직")

# Left Column: Inputs/Outputs, Right Column: Logic Process
il, it, iw, ih = draw_card(slide, Inches(0.8), top, w2, height, "공개 및 비밀 파라미터 구성")
tx = slide.shapes.add_textbox(il, it, iw, ih)
tf = tx.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "■ 회로에 사용되는 주요 파라미터 구조"
p.font.name = "Apple SD Gothic Neo"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = ACCENT_INDIGO

add_bullet(tf, "공개 파라미터 (Public Inputs): ", "검증자(DC)가 대조 가능한 값\n - 범위 임계치 (range.lower, range.upper)\n - 데이터 속성 식별자 (attribute)\n - 발급 타임스탬프 (timestamp)\n - 발급기관 등록 머클 루트 (merkle_root)")
add_bullet(tf, "비밀 파라미터 (Private Witnesses): ", "은닉해야 하는 값\n - 데이터 원문 수치 (value)\n - 발급기관 공개키 (di.public_key)\n - 머클 멤버십 증명 경로 및 비트 (merkle_path, merkle_bits)")

il, it, iw, ih = draw_card(slide, Inches(0.8) + w2 + g2, top, w2, height, "서킷 검증 연산 흐름")
tx = slide.shapes.add_textbox(il, it, iw, ih)
tf = tx.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "■ 3단계의 논리적 서킷 검증"
p.font.name = "Apple SD Gothic Neo"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = ACCENT_TEAL

add_bullet(tf, "1. 수치 범위 검증 (Range Test): ", "원문 value가 공개된 range.lower와 range.upper 사이에 존재하는지 수학적 부등식 연산 수행")
add_bullet(tf, "2. 발급서명 검증 (Signature Verification): ", "발급된 데이터가 신뢰할 수 있는 기관에 의해 서명되었는지 EdDSA 서명 및 MiMC 해시 알고리즘을 통해 서킷 내에서 검증")
add_bullet(tf, "3. 기관 유효성 검증 (Membership Proof): ", "비밀 입력인 di.public_key를 머클 경로를 따라 해싱하여 공개된 merkle_root와 일치하는지 검증 ($O(\log N)$ 복잡도로 효율성 확보)")


# ==========================================================
# SLIDE 9: 다중 쿼리 및 소유권 증명 (Composite Query & Ownership Proof)
# ==========================================================
slide = prs.slides.add_slide(slide_layout)
apply_slide_background(slide)
add_header(slide, "03. Proposed System", "다중 쿼리 처리와 소유권 증명 프로토콜")

il, it, iw, ih = draw_card(slide, Inches(0.8), top, w2, height, "소유권 증명 문제 (Ownership Proof Problem)")
tx = slide.shapes.add_textbox(il, it, iw, ih)
tf = tx.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "■ 복합 조건 검색 시 데이터 연계 입증 필요성"
p.font.name = "Apple SD Gothic Neo"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = ACCENT_ORANGE

add_bullet(tf, "프라이버시와 무결성의 충돌: ", "검색 결과 반환 시, 환자 식별자(DP_ID)는 유출을 방지하기 위해 완전 은닉되어야 함")
add_bullet(tf, "위조 및 조작 공격 가능성: ", "서버(SS)가 악의적으로 서로 다른 환자의 증명들을 무단 결합하여 다중 조건 만족 결과로 둔갑시키는 기만 행위 차단 필요")
add_bullet(tf, "보안 요구사항: ", "결과로 반환된 다수의 증명들이 실제로 '동일한 소유자(DP)'의 것임을 프라이버시 노출 없이 입증해야 함")

il, it, iw, ih = draw_card(slide, Inches(0.8) + w2 + g2, top, w2, height, "해결 방안: 단일 그룹 커밋 (Single Group Commitment)")
tx = slide.shapes.add_textbox(il, it, iw, ih)
tf = tx.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "■ 임의 난수를 이용한 소유권 수학적 결착"
p.font.name = "Apple SD Gothic Neo"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = ACCENT_TEAL

add_bullet(tf, "1. nonce 전달: ", "수요자(DC)가 다중 쿼리 전송 시 임의의 일회성 난수(nonce)를 동봉")
add_bullet(tf, "2. 증명 그룹화: ", "SS는 동일 DP_ID를 기반으로 쿼리에 부합하는 증명 식별자들($p_1, \dots, p_n$)을 필터링 및 조인")
add_bullet(tf, "3. prf_hash 계산: ", "그룹 내 증명 식별자들의 총합을 구하여 해싱 : $prf\\_hash = Hash(p_1 + \\dots + p_n)$")
add_bullet(tf, "4. 단일 커밋 도출: ", "최종 커밋 $c = Hash(prf\\_hash + DP\\_ID + nonce)$ 생성 후, DP_ID 대신 커밋 $c$를 결과로 반환")


# ==========================================================
# SLIDE 10: 소유권 증명 회로 및 최적화 (Ownership Proof Circuit & Optimization)
# ==========================================================
slide = prs.slides.add_slide(slide_layout)
apply_slide_background(slide)
add_header(slide, "03. Proposed System", "소유권 증명 회로 설계 및 실시간 최적화")

il, it, iw, ih = draw_card(slide, Inches(0.8), top, w2, height, "소유권 증명 산술 회로 구조")
tx = slide.shapes.add_textbox(il, it, iw, ih)
tf = tx.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "■ 동일인 및 데이터 관계 입증 로직"
p.font.name = "Apple SD Gothic Neo"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = ACCENT_INDIGO

add_bullet(tf, "공개 입력 (Public Inputs): ", " - 계산된 커밋 값 ($c$)\n - 결과 증명 식별자 리스트 ($p_1, \dots, p_n$)\n - 수요자가 전달한 난수 ($nonce$)")
add_bullet(tf, "비밀 입력 (Private Witness): ", " - 소유자의 고유 식별자 ($DP\\_ID$)")
add_bullet(tf, "회로 검증식: ", "비밀 입력된 $DP\\_ID$와 공개 입력들을 사용해 회로 내에서 커밋 $c$의 유효성 검사\n - $c == MiMCHash(MiMCHash(\\sum p_i) + DP\\_ID + nonce)$\n - 일치할 경우, 반환된 모든 증명이 해당 DP_ID의 것임을 인증")

il, it, iw, ih = draw_card(slide, Inches(0.8) + w2 + g2, top, w2, height, "실시간 일괄 생성 최적화 (On-the-fly Batch Proof)")
tx = slide.shapes.add_textbox(il, it, iw, ih)
tf = tx.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "■ 비대화형(Non-interactive) 1 round-trip 설계"
p.font.name = "Apple SD Gothic Neo"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = ACCENT_TEAL

add_bullet(tf, "초경량 서킷 설계: ", "소유권 회로는 MiMC 해시 연산만으로 구성되어 증명 생성 시간이 수 ms 수준으로 극도로 가벼움")
add_bullet(tf, "선제적 일괄 증명 동봉: ", "대화형 방식(On-demand 요청)을 탈피하여, SS가 검색 결과를 반환하는 시점에 실시간으로 소유권 증명을 선제 생성하여 전달")
add_bullet(tf, "통신 오버헤드 최소화: ", "추가 네트워크 라운드 트립을 제거하여 단 1회의 요청/응답(1 Round-trip)으로 완벽한 실시간 검색 달성")


# ==========================================================
# SLIDE 11: 실험 및 평가 - 마이크로 벤치마크 (Evaluation - Micro-benchmarks)
# ==========================================================
slide = prs.slides.add_slide(slide_layout)
apply_slide_background(slide)
add_header(slide, "04. Evaluation", "실험 및 평가: 마이크로 벤치마크 분석")

# Left: Range Proof Benchmark Table & Details, Right: Ownership Proof Table & Details
w_sub = Inches(5.65)
h_sub2 = Inches(2.2)

# Range Proof Card (Top Left)
il, it, iw, ih = draw_card(slide, Inches(0.8), Inches(1.6), w_sub, h_sub2, "1. 범위 증명 마이크로 벤치마크 (Range Proof)")
rows, cols = 3, 5
left_t, top_t, width_t, height_t = il + Inches(0.1), it + Inches(0.5), iw - Inches(0.2), Inches(0.9)
table_shape = slide.shapes.add_table(rows, cols, left_t, top_t, width_t, height_t)
table = table_shape.table
table.columns[0].width = Inches(1.8)
for i in range(1, 5):
    table.columns[i].width = Inches(0.8)

headers = ["DI 개수", "128", "256", "512", "1024"]
row1 = ["증명생성시간 (ms)", "75", "80", "83", "86"]
row2 = ["증명 크기 (Bytes)", "164", "164", "164", "164"]

for c_idx in range(5):
    cell = table.cell(0, c_idx)
    cell.text = headers[c_idx]
    cell.fill.solid()
    cell.fill.fore_color.rgb = CARD_BG
    p = cell.text_frame.paragraphs[0]
    p.font.size = Pt(10)
    p.font.name = "Apple SD Gothic Neo"
    p.font.bold = True
    p.font.color.rgb = ACCENT_TEAL
    p.alignment = PP_ALIGN.CENTER
    
    cell1 = table.cell(1, c_idx)
    cell1.text = row1[c_idx]
    cell1.fill.solid()
    cell1.fill.fore_color.rgb = CARD_BG
    p1 = cell1.text_frame.paragraphs[0]
    p1.font.size = Pt(9)
    p1.font.name = "Apple SD Gothic Neo"
    p1.font.color.rgb = TEXT_PRIMARY
    p1.alignment = PP_ALIGN.CENTER
    
    cell2 = table.cell(2, c_idx)
    cell2.text = row2[c_idx]
    cell2.fill.solid()
    cell2.fill.fore_color.rgb = CARD_BG
    p2 = cell2.text_frame.paragraphs[0]
    p2.font.size = Pt(9)
    p2.font.name = "Apple SD Gothic Neo"
    p2.font.color.rgb = TEXT_PRIMARY
    p2.alignment = PP_ALIGN.CENTER

tx_box = slide.shapes.add_textbox(il, it + Inches(1.4), iw, Inches(0.7))
tf_tx = tx_box.text_frame
tf_tx.word_wrap = True
p = tf_tx.paragraphs[0]
p.text = "• DI 수가 8배 늘어날 때, 연산 지연은 단 11ms (15%) 증가에 그침 (머클 증명 O(log N) 확장성 증명)\n• 증명 크기는 Groth16 프로토콜 특성상 164바이트로 항시 일정"
p.font.size = Pt(10)
p.font.name = "Apple SD Gothic Neo"
p.font.color.rgb = TEXT_SECONDARY

# Ownership Proof Card (Bottom Left)
il2, it2, iw2, ih2 = draw_card(slide, Inches(0.8), Inches(4.2), w_sub, h_sub2, "2. 소유권 증명 마이크로 벤치마크 (Ownership Proof)")
table_shape2 = slide.shapes.add_table(3, 5, il2 + Inches(0.1), it2 + Inches(0.5), iw2 - Inches(0.2), Inches(0.9))
table2 = table_shape2.table
table2.columns[0].width = Inches(1.8)
for i in range(1, 5):
    table2.columns[i].width = Inches(0.8)

headers2 = ["서브쿼리 개수", "2", "3", "4", "5"]
row1_2 = ["증명생성시간 (ms)", "10", "11", "12", "16"]
row2_2 = ["증명 크기 (Bytes)", "164", "164", "164", "164"]

for c_idx in range(5):
    cell = table2.cell(0, c_idx)
    cell.text = headers2[c_idx]
    cell.fill.solid()
    cell.fill.fore_color.rgb = CARD_BG
    p = cell.text_frame.paragraphs[0]
    p.font.size = Pt(10)
    p.font.name = "Apple SD Gothic Neo"
    p.font.bold = True
    p.font.color.rgb = ACCENT_TEAL
    p.alignment = PP_ALIGN.CENTER
    
    cell1 = table2.cell(1, c_idx)
    cell1.text = row1_2[c_idx]
    cell1.fill.solid()
    cell1.fill.fore_color.rgb = CARD_BG
    p1 = cell1.text_frame.paragraphs[0]
    p1.font.size = Pt(9)
    p1.font.name = "Apple SD Gothic Neo"
    p1.font.color.rgb = TEXT_PRIMARY
    p1.alignment = PP_ALIGN.CENTER
    
    cell2 = table2.cell(2, c_idx)
    cell2.text = row2_2[c_idx]
    cell2.fill.solid()
    cell2.fill.fore_color.rgb = CARD_BG
    p2 = cell2.text_frame.paragraphs[0]
    p2.font.size = Pt(9)
    p2.font.name = "Apple SD Gothic Neo"
    p2.font.color.rgb = TEXT_PRIMARY
    p2.alignment = PP_ALIGN.CENTER

tx_box2 = slide.shapes.add_textbox(il2, it2 + Inches(1.4), iw2, Inches(0.7))
tf_tx2 = tx_box2.text_frame
tf_tx2.word_wrap = True
p = tf_tx2.paragraphs[0]
p.text = "• 서브 쿼리 개수 증가에 따라 해시 연산량이 증가하나 최대 16ms로 실시간 생성에 부담 없음\n• 30건 일괄 생성 시에도 500ms 미만으로 연산 완료"
p.font.size = Pt(10)
p.font.name = "Apple SD Gothic Neo"
p.font.color.rgb = TEXT_SECONDARY

# Right Side Column: Composite Query & Storage Overhead (Vertical Card)
il_r, it_r, iw_r, ih_r = draw_card(slide, Inches(0.8) + w_sub + g, Inches(1.6), w_sub, Inches(4.8), "다중 쿼리 응답 분석 및 용량 오버헤드")
tx_r = slide.shapes.add_textbox(il_r, it_r, iw_r, ih_r)
tf_r = tx_r.text_frame
tf_r.word_wrap = True
p = tf_r.paragraphs[0]
p.text = "■ 다중 쿼리 (Composite Query) 전체 처리 지연"
p.font.name = "Apple SD Gothic Neo"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = ACCENT_INDIGO

add_bullet(tf_r, "DB 조회의 큰 영향: ", "서브쿼리 2개(3.7초) -> 5개(9.0초)로 증가 시 전체 시간 중 DB Fetch(조인 및 필터링)가 85~95% 점유\n - 순수 소유권 증명 시간은 ~450ms로 거의 일정")
add_bullet(tf_r, "벤치마크 격차 분석: ", "캐시 미스 및 컨텍스트 스위칭 등의 I/O 오버헤드로 인해 단일 증명 반복 대비 전체 지연이 소폭 증가하였으나 실용적 범위 이내")

p_st = tf_r.add_paragraph()
p_st.text = "■ 스토리지 및 네트워크 전송 용량 분석"
p_st.font.name = "Apple SD Gothic Neo"
p_st.font.size = Pt(13)
p_st.font.bold = True
p_st.font.color.rgb = ACCENT_INDIGO
p_st.space_before = Pt(10)

add_bullet(tf_r, "스토리지 총 용량: ", "100만 건 데이터 기준, 공개 파라미터를 포함한 ZKRP 전체 스토리지 크기는 약 630MB\n - 순수 증명만 저장 시 164MB로 감소 가능")
add_bullet(tf_r, "네트워크 저부하: ", "페이지 크기 30일 때 반환되는 소유권 증명들의 크기는 약 4.9KB 수준으로 일반 이미지 1장 이하의 네트워크 데이터 전송량 발생")


# ==========================================================
# SLIDE 12: 실험 및 평가 - 설계 방식 비교 (Evaluation - Comparison)
# ==========================================================
slide = prs.slides.add_slide(slide_layout)
apply_slide_background(slide)
add_header(slide, "04. Evaluation", "실험 및 평가: 설계 방식(Proposed vs Baseline) 비교")

# 2 Columns (Left: Table of Results, Right: Detailed Analysis)
il, it, iw, ih = draw_card(slide, Inches(0.8), top, w2, height, "성능 비교 측정 데이터 (단일 쿼리 100회 평균, DI 128, Page 30)")

# Create Table for Comparison
rows, cols = 8, 4
left_t, top_t, width_t, height_t = il + Inches(0.1), it + Inches(0.2), iw - Inches(0.2), Inches(3.8)
table_shape = slide.shapes.add_table(rows, cols, left_t, top_t, width_t, height_t)
table = table_shape.table
table.columns[0].width = Inches(1.8)
for i in range(1, 4):
    table.columns[i].width = Inches(1.2)

headers = ["지표 (ms)", "Baseline (Raw)", "Baseline (Enc)", "제안 시스템 (Proposed)"]
col_data = [
    ["DB Fetch (평균)", "97", "269", "31"],
    ["Overall (평균)", "2242", "2405", "315"],
    ["편차 (StdDev)", "209", "203", "72"],
    ["p50 (중앙값)", "2245", "2411", "292"],
    ["p90 (백분위)", "2326", "2444", "360"],
    ["p95 (백분위)", "2370", "2485", "372"],
    ["p99 (꼬리지연)", "2460", "2806", "390"]
]

for col_idx in range(4):
    cell = table.cell(0, col_idx)
    cell.text = headers[col_idx]
    cell.fill.solid()
    cell.fill.fore_color.rgb = CARD_BG
    p = cell.text_frame.paragraphs[0]
    p.font.size = Pt(10)
    p.font.name = "Apple SD Gothic Neo"
    p.font.bold = True
    p.font.color.rgb = ACCENT_TEAL
    p.alignment = PP_ALIGN.CENTER

for row_idx in range(7):
    for col_idx in range(4):
        cell = table.cell(row_idx + 1, col_idx)
        cell.text = col_data[row_idx][col_idx]
        cell.fill.solid()
        if col_idx == 3:
            cell.fill.fore_color.rgb = RGBColor(38, 55, 78)
        else:
            cell.fill.fore_color.rgb = CARD_BG
            
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(9)
        p.font.name = "Apple SD Gothic Neo"
        if col_idx == 3:
            p.font.bold = True
            p.font.color.rgb = ACCENT_TEAL
        else:
            p.font.color.rgb = TEXT_PRIMARY
        p.alignment = PP_ALIGN.CENTER

# Right card: Key Insights
il2, it2, iw2, ih2 = draw_card(slide, Inches(0.8) + w2 + g2, top, w2, height, "설계 방식별 비교 결과 심층 분석")
tx2 = slide.shapes.add_textbox(il2, it2, iw2, ih2)
tf2 = tx2.text_frame
tf2.word_wrap = True
p = tf2.paragraphs[0]
p.text = "■ 7~8배 빠른 전체 검색 응답 속도 달성"
p.font.name = "Apple SD Gothic Neo"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = ACCENT_TEAL

add_bullet(tf2, "ZKP 실시간 오버헤드 완벽 제거: ", "Baseline 모델들은 쿼리 결과를 확인한 후 실시간으로 범위 증명을 생성하므로 약 2.2초의 생성 병목 발생\n - 제안 모델은 사전 증명 인덱스 검색(31ms)만 수행하므로 전체 315ms 달성")
add_bullet(tf2, "AES-CTR 암호화 복호화 지연: ", "Baseline(Enc)의 경우 모든 튜플을 복호화하며 스캔해야 하므로 DB Fetch에만 269ms가 소요되어 성능 저하")
add_bullet(tf2, "성능 편차 및 안정성 확보: ", "영지식 연산이 빠진 제안 모델은 표준편차가 72ms로 가장 낮고, 꼬리 지연(p99) 역시 390ms로 극단적 상황에서도 매우 안정적임")


# ==========================================================
# SLIDE 13: 실험 및 평가 - 확장성 분석 (Evaluation - Scalability)
# ==========================================================
slide = prs.slides.add_slide(slide_layout)
apply_slide_background(slide)
add_header(slide, "04. Evaluation", "실험 및 평가: 확장성 분석 (DI 수 & 페이지 크기)")

il, it, iw, ih = draw_card(slide, Inches(0.8), top, w2, height, "1. 발급자(DI) 수 증가에 따른 응답 지연 추이")
tx = slide.shapes.add_textbox(il, it, iw, ih)
tf = tx.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "■ DI 수 확장 시 응답 속도 변화 (128 -> 1024)"
p.font.name = "Apple SD Gothic Neo"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = ACCENT_INDIGO

add_bullet(tf, "Proposed (제안): ", "128 DI(315ms) -> 1024 DI(300ms) 대의 안정적인 속도를 유지하며 확장에 대한 무영향성 입증")
add_bullet(tf, "Baseline (Raw/Enc): ", "128 DI(2242ms) -> 1024 DI(2602ms) 수준으로 증가. 머클 트리 깊이가 깊어짐에 따라 실시간 범위 증명 서킷의 연산 비용이 함께 누적됨")
add_bullet(tf, "확장성 평가: ", "제안 모델은 발급자 수에 종속되지 않는 독립적인 O(1) 수준의 검색 질의 속도를 실현함")

il2, it2, iw2, ih2 = draw_card(slide, Inches(0.8) + w2 + g2, top, w2, height, "2. 페이지 크기(Page Size) 증가에 따른 응답 지연 추이")
tx2 = slide.shapes.add_textbox(il2, it2, iw2, ih2)
tf2 = tx2.text_frame
tf2.word_wrap = True
p = tf2.paragraphs[0]
p.text = "■ 반환 건수 증가에 따른 선형 vs 완만 성장 구도 (30 -> 120)"
p.font.name = "Apple SD Gothic Neo"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = ACCENT_TEAL

add_bullet(tf2, "Baseline의 한계 (O(K) 선형 증가): ", "페이지 크기가 30에서 120으로 증가함에 따라 실시간 ZKP 생성 개수가 4배 증가하여 전체 시간이 약 2.6초에서 8.0초로 선형 증가")
add_bullet(tf2, "Proposed의 이점 (O(1) 수렴): ", "페이지 크기가 30에서 120으로 증가하더라도 응답 시간은 276ms에서 402ms로 극히 완만하게 증가(DB 조회 오버헤드만 추가)")
add_bullet(tf2, "대용량 서비스 적용 가능성: ", "사용자가 대량의 데이터를 검색하는 상용 거래 플랫폼 환경에 제안 모델이 필수적임을 수학적으로 입증")


# ==========================================================
# SLIDE 14: 보안성 고찰 및 위협 모델 (Security Analysis & Trade-off)
# ==========================================================
slide = prs.slides.add_slide(slide_layout)
apply_slide_background(slide)
add_header(slide, "04. Evaluation", "보안성 고찰: 추론 공격 시뮬레이션 및 대응 전략")

il, it, iw, ih = draw_card(slide, Inches(0.8), top, w2, height, "추론 공격 시뮬레이션 (Inference Attack)")
tx = slide.shapes.add_textbox(il, it, iw, ih)
tf = tx.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "■ 서버 장악 시 소유권 맵핑 테이블(OWNERSHIP) 유출 위협"
p.font.name = "Apple SD Gothic Neo"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = ACCENT_ORANGE

add_bullet(tf, "공격자 획득 지식: ", "공격자는 개별 DP_ID가 소유한 모든 범위 증명 리스트를 분류하고, 임상참고치의 경계를 결합하여 데이터 분포 확인 가능")
add_bullet(tf, "개인 건강 추론 예시: ", "100만 건 데이터 중 특정 환자의 공복혈당 ZKRP 37개 분포를 추출한 결과, 당뇨 진단 범주(125 이상)에 해당하는 기록이 14건 발견되어 병력 추론 가능")
add_bullet(tf, "현실적 위협: ", "정확한 수치(평문)는 가려지나, 사전 정의 구간 기반 특성으로 인해 유병 여부 등의 민감한 프라이버시가 상당 부분 복원될 위험 잔존")

il2, it2, iw2, ih2 = draw_card(slide, Inches(0.8) + w2 + g2, top, w2, height, "보안 완화 전략 및 시스템적 트레이드오프")
tx2 = slide.shapes.add_textbox(il2, it2, iw2, ih2)
tf2 = tx2.text_frame
tf2.word_wrap = True
p = tf2.paragraphs[0]
p.text = "■ 프라이버시 방어책의 기술적 절충안"
p.font.name = "Apple SD Gothic Neo"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = ACCENT_TEAL

add_bullet(tf2, "1. 더미 데이터 주입 (Dummy Injection): ", "실제 데이터 외에 수학적으로 참이지만 거짓인 'Dummy ZKRP'를 1:1로 생성해 노이즈를 섞음. 공격자가 실데이터를 식별 불가능하게 교란하나 구매 결정 후 진짜를 필터링하는 로직 필요")
add_bullet(tf2, "2. 비공모 다중 서버 비밀 분산: ", "OWNERSHIP 소유자 매핑 관계를 단일 서버가 아닌 독립된 다중 서버에 MPC(다자간 연산) 기반으로 분산하여 저장 및 연산. 단일 서버가 장악되더라도 프라이버시가 완벽히 보장되나 서버 연산 오버헤드 증가")
add_bullet(tf2, "결과적 타협: ", "본 연구는 실시간성 확보를 위해 반신뢰(Semi-honest) 모델을 기본 상정하여 전략적 합의를 선택함")


# ==========================================================
# SLIDE 15: 결론 및 향후 과제 (Conclusion & Future Work)
# ==========================================================
slide = prs.slides.add_slide(slide_layout)
apply_slide_background(slide)
add_header(slide, "05. Conclusion", "결론 및 향후 과제")

il, it, iw, ih = draw_card(slide, Inches(0.8), top, w2, height, "연구 요약 및 학술적/실용적 의의")
tx = slide.shapes.add_textbox(il, it, iw, ih)
tf = tx.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "■ 프라이버시와 검색 가능성의 성공적 조율"
p.font.name = "Apple SD Gothic Neo"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = ACCENT_TEAL

add_bullet(tf, "근본적 모순 해결: ", "데이터 원문을 서버에 저장하지 않는 '범위 증명 검색' 아키텍처를 도입하여 데이터 최소화와 프라이버시 보호 실현")
add_bullet(tf, "실시간 연산 병목 제거: ", "의료 참고치를 이용한 오프라인 사전 ZKRP 생성을 통해 쿼리 응답 시간을 약 7~8배 단축하고, 수백 ms 이내에 동적 데이터 검색 달성")
add_bullet(tf, "소유권 수학적 검증: ", "단일 그룹 커밋 및 경량 소유권 회로를 도입하여 다중 쿼리 결과를 1-round-trip 비대화형으로 검증 완료하여 실용성 입증")

il2, it2, iw2, ih2 = draw_card(slide, Inches(0.8) + w2 + g2, top, w2, height, "연구의 한계점 및 향후 연구 과제")
tx2 = slide.shapes.add_textbox(il2, it2, iw2, ih2)
tf2 = tx2.text_frame
tf2.word_wrap = True
p = tf2.paragraphs[0]
p.text = "■ 고도화 및 범용성 확보를 위한 향후 과제"
p.font.name = "Apple SD Gothic Neo"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = ACCENT_INDIGO

add_bullet(tf2, "한계 1 (소유권 유출 위협): ", "실시간 그룹화를 위한 OWNERSHIP 테이블 관리로 인해 서버 장악 시 데이터 분포 노출 위험 잔존")
add_bullet(tf2, "한계 2 (반신뢰 서버 가정): ", "결과의 위조/조작(Soundness)은 보장하지만 악의적 서버의 고의 누락(Completeness)은 방어하지 못함")
add_bullet(tf2, "향후 연구 과제 1: ", "검색의 실시간성을 훼손하지 않으면서 결과의 누락 없음(Completeness)을 수학적으로 검증하는 완결성 프로토콜 개발")
add_bullet(tf2, "향후 연구 과제 2: ", "이산적 구간(Pre-defined Range) 중심의 설계에서 나아가 연속 범위 검색 및 비수치형(텍스트 등) 데이터에 대한 범용 영지식 검색 연구")


# 4. Save Presentation
output_path = "/Users/jinwoo/Desktop/Master-s-Degree-Paper/presentation.pptx"
prs.save(output_path)
print(f"Presentation saved successfully to {output_path}")
